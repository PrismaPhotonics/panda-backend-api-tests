"""
Generate PowerPoint Presentation with Code Links
================================================

Creates a professional presentation with:
1. Direct links to development code
2. QA tests without specs summary
3. Detailed CSV of all affected tests

Usage:
    python scripts/generate_presentation_with_links.py

Requirements:
    pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import csv
from pathlib import Path
from datetime import datetime

# =============================================================================
# Configuration
# =============================================================================

# Link Mode: Choose how to open code files
LINK_MODE = "local"  # Options: "local" (Cursor/VS Code) or "web" (GitHub/Bitbucket)

# Local Configuration (for Cursor/VS Code links)
LOCAL_PROJECT_PATH = r"C:\Projects\focus_server_automation"

# Web Configuration (for GitHub/Bitbucket links)
REPO_BASE_URL = "https://github.com/your-org/focus_server_automation/blob/main/"
# Or for Bitbucket: "https://bitbucket.company.com/projects/PZ/repos/focus-server/browse/"

# Output paths
OUTPUT_DIR = Path("documentation/analysis")
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

DECK_PATH = OUTPUT_DIR / "Automation_Specs_Gap_Review_LOCAL_LINKS.pptx"
CSV_PATH = OUTPUT_DIR / "tests_without_specs_LOCAL_LINKS.csv"


# =============================================================================
# Helper Functions
# =============================================================================

def create_code_link(file_path: str, line_range: str = None) -> str:
    """
    Create a link to code file.
    
    Args:
        file_path: Relative path from project root (e.g., "src/utils/validators.py")
        line_range: Line number or range (e.g., "395" or "390-460")
    
    Returns:
        URL string (local vscode:// or web GitHub/Bitbucket link)
    """
    if LINK_MODE == "local":
        # Convert Windows path to forward slashes for VS Code URI
        full_path = Path(LOCAL_PROJECT_PATH) / file_path
        path_str = str(full_path).replace("\\", "/")
        
        if line_range:
            # Extract first line number from range (e.g., "390-460" -> "390")
            line_num = line_range.split("-")[0] if "-" in line_range else line_range
            return f"vscode://file/{path_str}:{line_num}"
        else:
            return f"vscode://file/{path_str}"
    
    else:  # web mode
        if line_range:
            # GitHub format: #L390-L460
            if "-" in line_range:
                start, end = line_range.split("-")
                return f"{REPO_BASE_URL}{file_path}#L{start}-L{end}"
            else:
                return f"{REPO_BASE_URL}{file_path}#L{line_range}"
        else:
            return f"{REPO_BASE_URL}{file_path}"


def create_presentation():
    """Create PowerPoint presentation with code links and test details."""
    
    # Create new presentation
    prs = Presentation()
    
    # ========================================================================
    # Slide 1: Direct Dev Code Links
    # ========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only layout
    title = slide.shapes.title
    title.text = "Direct Dev Code Links (Evidence)"
    title.text_frame.paragraphs[0].font.name = "Arial"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    
    # Add text box with links
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9.1), Inches(4.8))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    
    code_links = [
        ("[P0] #1: ROI change limit - hardcoded 50%", 
         create_code_link("src/utils/validators.py", "390-460")),
        
        ("[P0] #2: Performance assertions disabled (P95/P99)", 
         create_code_link("tests/integration/performance/test_performance_high_priority.py", "146-170")),
        
        ("[P1] #3: NFFT validation too permissive", 
         create_code_link("src/utils/validators.py", "194-227")),
        
        ("[P1] #4: Frequency range - no absolute max/min", 
         create_code_link("src/models/focus_server_models.py", "46-57")),
        
        ("[P2] #5: Sensor range - no min/max ROI size", 
         create_code_link("src/utils/validators.py", "116-151")),
        
        ("[P2] #6: Polling helper - hardcoded timeouts", 
         create_code_link("src/utils/helpers.py", "474-504")),
        
        ("[P2] #7: Default payloads mismatch config", 
         create_code_link("src/utils/helpers.py", "507-532")),
        
        ("[P3] Config validation tests with TODO/no assertions", 
         create_code_link("tests/integration/api/test_config_validation_high_priority.py", "475-520")),
        
        ("[P3] MongoDB outage resilience (behavior unclear)", 
         create_code_link("tests/integration/performance/test_mongodb_outage_resilience.py")),
    ]
    
    for i, (label, url) in enumerate(code_links):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        
        # Label part
        r1 = p.add_run()
        r1.text = f"{label}  "
        r1.font.name = "Arial"
        r1.font.size = Pt(20)
        
        # Link part
        r2 = p.add_run()
        r2.text = "[view code →]"
        r2.font.name = "Arial"
        r2.font.size = Pt(18)
        r2.font.color.rgb = RGBColor(0, 102, 204)  # Blue
        r2.hyperlink.address = url
    
    # ========================================================================
    # Slide 2: QA Tests Without Specs – Summary (WITH LINKS)
    # ========================================================================
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Title only layout
    title2 = slide2.shapes.title
    title2.text = "QA Tests Without Specs – Evidence (Summary)"
    title2.text_frame.paragraphs[0].font.name = "Arial"
    title2.text_frame.paragraphs[0].font.size = Pt(32)
    title2.text_frame.paragraphs[0].font.bold = True
    
    # Add text box with links
    tb2 = slide2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9.1), Inches(4.8))
    tf2 = tb2.text_frame
    tf2.clear()
    tf2.word_wrap = True
    
    tests_with_links = [
        ("test_p95_p99_latency_post_config - P95/P99 thresholds missing",
         create_code_link("tests/integration/performance/test_performance_high_priority.py", "146-170"),
         RGBColor(220, 53, 69)),  # Red
        
        ("test_get_channels_endpoint_response_time - SLA missing",
         create_code_link("tests/integration/api/test_api_endpoints_high_priority.py", "135-147"),
         RGBColor(220, 53, 69)),  # Red
        
        ("test_mongodb_scale_down_outage - expected HTTP code/recovery unknown",
         create_code_link("tests/integration/performance/test_mongodb_outage_resilience.py"),
         RGBColor(253, 126, 20)),  # Orange
        
        ("test_concurrent_task_max_limit - system limit undefined",
         create_code_link("tests/integration/performance/test_concurrency_limits.py"),
         RGBColor(253, 126, 20)),  # Orange
        
        ("test_waterfall_poll_200_until_data - max wait/alert/timeout undefined",
         create_code_link("tests/integration/api/test_waterfall_polling_behavior.py"),
         RGBColor(255, 193, 7)),  # Yellow
        
        ("test_frequency_range_equal_min_max - behavior for min==max undefined",
         create_code_link("tests/integration/api/test_config_validation_high_priority.py", "475-495"),
         RGBColor(255, 193, 7)),  # Yellow
        
        ("test_channel_range_equal_min_max - behavior for min==max undefined",
         create_code_link("tests/integration/api/test_config_validation_high_priority.py", "506-520"),
         RGBColor(255, 193, 7)),  # Yellow
        
        ("test_config_with_invalid_sensor_range - ROI min/max size undefined",
         create_code_link("tests/integration/api/test_config_validation_boundaries.py"),
         RGBColor(255, 193, 7)),  # Yellow
        
        ("test_nfft_values_enforcement - valid values list not enforced",
         create_code_link("tests/integration/api/test_signal_params_validation.py"),
         RGBColor(255, 193, 7)),  # Yellow
        
        ("test_error_schema_consistency - API error contract undefined",
         create_code_link("tests/integration/api/test_error_schema_contract.py"),
         RGBColor(255, 193, 7)),  # Yellow
    ]
    
    for i, (label, url, color) in enumerate(tests_with_links):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.level = 0
        
        # Label part with color
        r1 = p.add_run()
        r1.text = f"{label}  "
        r1.font.name = "Arial"
        r1.font.size = Pt(16)
        r1.font.color.rgb = color
        
        # Link part
        r2 = p.add_run()
        r2.text = "[view code →]"
        r2.font.name = "Arial"
        r2.font.size = Pt(14)
        r2.font.color.rgb = RGBColor(0, 102, 204)  # Blue
        r2.hyperlink.address = url
    
    # Save presentation
    prs.save(DECK_PATH)
    print(f"[OK] PowerPoint presentation created: {DECK_PATH}")
    
    return DECK_PATH


def create_detailed_csv():
    """Create detailed CSV of all tests without specs."""
    
    csv_headers = [
        "Test Name / ID",
        "Path",
        "What It Checks",
        "Missing Spec / Decision Needed",
        "Impact (Why It Matters)",
        "Linked Code (Ref)"
    ]
    
    tests_rows = [
        [
            "test_p95_p99_latency_post_config (PZ-13770)",
            "tests/integration/performance/test_performance_high_priority.py:146-170",
            "Measures P95/P99 for POST /config",
            "Define P95/P99 thresholds + timeout + acceptable error rate",
            "Cannot enforce performance regression gates without SLA",
            "src/utils/helpers.py (latency harness), API handlers for /config"
        ],
        [
            "test_get_channels_endpoint_response_time (PZ-13419.1)",
            "tests/integration/api/test_api_endpoints_high_priority.py:135-147",
            "Response time for GET /channels",
            "Max response time SLA for endpoint",
            "Arbitrary 1000ms today; could be too strict/lenient",
            "FastAPI route for /channels"
        ],
        [
            "test_mongodb_scale_down_outage_returns_503",
            "tests/integration/performance/test_mongodb_outage_resilience.py",
            "Behavior under MongoDB outage (status, recovery time)",
            "Expected HTTP code (e.g., 503), max recovery time, caching policy",
            "Outage handling unclear → flaky triage, undefined UX",
            "DB client layer, error handlers"
        ],
        [
            "test_concurrent_task_max_limit",
            "tests/integration/performance/test_concurrency_limits.py",
            "System behavior under N concurrent tasks",
            "Define max concurrent tasks and on-exceed behavior",
            "Unknown capacity → risk of overload or false failures",
            "Task orchestrator / queue config"
        ],
        [
            "test_waterfall_poll_200_until_data",
            "tests/integration/api/test_waterfall_polling_behavior.py",
            "Polling semantics when 200-no-data-yet",
            "Max wait before alert/timeout (live vs historic) and backoff",
            "Tests can hang or be flaky without a contract",
            "poll_until helper; /waterfall handler"
        ],
        [
            "test_frequency_range_equal_min_max (PZ-13876.1)",
            "tests/integration/api/test_config_validation_high_priority.py:475-495",
            "Edge case: frequency min==max",
            "Define validity and expected status code/message",
            "No assertion possible until contract defined",
            "FrequencyRange model; validators"
        ],
        [
            "test_channel_range_equal_min_max (PZ-13876.2)",
            "tests/integration/api/test_config_validation_high_priority.py:506-520",
            "Edge case: channel min==max",
            "Define validity (SingleChannel?) and expected behavior",
            "Cannot assert success/failure deterministically",
            "Channel range validation; config schema"
        ],
        [
            "test_config_with_invalid_sensor_range",
            "tests/integration/api/test_config_validation_boundaries.py",
            "Validate ROI sensor min/max against system constraints",
            "Define min/max ROI size (count) + total_sensors rules",
            "No guardrails → either false rejects or unsafe accepts",
            "validators.validate_sensor_range; settings.yaml constraints"
        ],
        [
            "test_nfft_values_enforcement",
            "tests/integration/api/test_signal_params_validation.py",
            "Accept only valid NFFT values",
            "Enforce allowed list (powers of 2) and max bound",
            "Invalid NFFT harms performance/memory characteristics",
            "validators.validate_nfft_value; settings.yaml:nfft.valid_values"
        ],
        [
            "test_error_schema_consistency",
            "tests/integration/api/test_error_schema_contract.py",
            "Standard error response mapping",
            "Define API error schema (code/message/detail/hint)",
            "Parsing brittle; automation must normalize errors",
            "Global exception handlers; Pydantic/FastAPI responses"
        ],
        [
            "validate_roi_change_safety (function)",
            "src/utils/validators.py:390-460",
            "Validates ROI changes are not too drastic",
            "Confirm max_change_percent (currently 50%), cooldown period",
            "6 tests depend on unconfirmed hardcoded value",
            "Used by ROI adjustment tests"
        ],
        [
            "validate_nfft_value (function)",
            "src/utils/validators.py:194-227",
            "Validates NFFT values",
            "Enforce valid_values list or keep warning-only behavior",
            "Code accepts any positive int; config has specific list",
            "Used by configuration validation"
        ],
        [
            "FrequencyRange (model)",
            "src/models/focus_server_models.py:46-57",
            "Pydantic model for frequency configuration",
            "Add absolute max/min frequency, minimum range span",
            "Currently accepts any positive values including extremes",
            "Used by all frequency-related tests"
        ],
        [
            "validate_sensor_range (function)",
            "src/utils/validators.py:116-151",
            "Validates sensor range configuration",
            "Define min/max ROI size (not just bounds)",
            "Could allow ROI with 1 sensor or all 2222 sensors",
            "Used by sensor configuration tests"
        ],
        [
            "poll_until (helper)",
            "src/utils/helpers.py:474-504",
            "Generic polling function",
            "Define timeout/interval for live vs historic scenarios",
            "Hardcoded 60s timeout, 1s interval - might not fit all cases",
            "Used throughout integration tests"
        ],
        [
            "generate_config_payload (helper)",
            "src/utils/helpers.py:507-532",
            "Generates test configuration payloads",
            "Align defaults with config/environments.yaml",
            "Code defaults don't match config file defaults",
            "Used by most integration tests"
        ],
    ]
    
    with open(CSV_PATH, "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow(csv_headers)
        writer.writerows(tests_rows)
    
    print(f"[OK] CSV file created: {CSV_PATH}")
    print(f"     Rows: {len(tests_rows)} tests/functions")
    
    return CSV_PATH


def main():
    """Main execution."""
    print("=" * 70)
    print("Generating Automation Specs Gap Review Materials")
    print("=" * 70)
    print()
    
    # Create presentation
    deck_path = create_presentation()
    
    # Create CSV
    csv_path = create_detailed_csv()
    
    print()
    print("=" * 70)
    print("Generation Complete!")
    print("=" * 70)
    print()
    print("Files created:")
    print(f"   1. PowerPoint: {deck_path}")
    print(f"   2. CSV:        {csv_path}")
    print()
    print("Next steps:")
    print("   1. Update REPO_BASE_URL in this script with your actual repo URL")
    print("   2. Open the PowerPoint file")
    print("   3. Test the code links (click [view code ->])")
    print("   4. Import CSV into Excel/Google Sheets for detailed review")
    print()
    print("To update links:")
    print("   - GitHub: https://github.com/your-org/repo/blob/main/")
    print("   - Bitbucket: https://bitbucket.company.com/projects/X/repos/Y/browse/")
    print()


if __name__ == "__main__":
    main()


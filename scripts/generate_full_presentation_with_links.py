#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Generate FULL 20-Slide PowerPoint with Code Links
==================================================

Creates complete presentation with code links in ALL relevant slides.
Based on: SLIDES_CONTENT_FOR_PRESENTATION.md

Usage:
    python scripts/generate_full_presentation_with_links.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

# Configuration
LOCAL_PROJECT_PATH = r"C:\Projects\focus_server_automation"
OUTPUT_DIR = Path("documentation/analysis")
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
DECK_PATH = OUTPUT_DIR / "Automation_Specs_Gap_Review_FULL_WITH_LINKS.pptx"


def create_code_link(file_path: str, line_range: str = None) -> str:
    """Create VS Code link to file."""
    full_path = Path(LOCAL_PROJECT_PATH) / file_path
    path_str = str(full_path).replace("\\", "/")
    
    if line_range:
        line_num = line_range.split("-")[0] if "-" in line_range else line_range
        return f"vscode://file/{path_str}:{line_num}"
    return f"vscode://file/{path_str}"


def add_title_slide(prs):
    """Slide 1: Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "AUTOMATION SPECS GAP REVIEW"
    p.font.name = "Arial"
    p.font.size = Pt(40)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    tb2 = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Focus Server Test Suite - Missing Specifications\nOctober 22, 2025"
    p2.font.name = "Arial"
    p2.font.size = Pt(24)
    p2.alignment = PP_ALIGN.CENTER


def add_executive_summary(prs):
    """Slide 2: Executive Summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Executive Summary"
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    content = [
        ("The Problem:", "", None),
        ("• We have 190+ automated tests", "", None),
        ("• But many lack clear PASS/FAIL criteria", "", None),
        ("• Due to missing specifications", "", None),
        ("", "", None),
        ("The Numbers:", "", None),
        ("• 82+ tests directly affected", "", None),
        ("• 50+ hardcoded values without confirmation", "", None),
        ("• 11 TODO comments waiting for specs", "", None),
        ("• 28 performance tests with disabled assertions", "", None),
    ]
    
    for i, (text, link, color) in enumerate(content):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if text:
            r = p.add_run()
            r.text = text
            r.font.name = "Arial"
            r.font.size = Pt(20) if ":" in text else Pt(18)
            r.font.bold = ":" in text


def add_code_evidence_1(prs):
    """Slide 3: Performance Tests Disabled - WITH LINK"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Evidence from Code #1: Performance Tests Disabled"
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(2))
    tf = tb.text_frame
    
    # Code example
    p1 = tf.paragraphs[0]
    r = p1.add_run()
    r.text = "# tests/integration/performance/test_performance_high_priority.py:157\n\n"
    r.font.name = "Courier New"
    r.font.size = Pt(16)
    r.font.color.rgb = RGBColor(100, 100, 100)
    
    r2 = p1.add_run()
    r2.text = "# TODO: Uncomment after specs meeting\n# assert p95 < THRESHOLD_P95_MS   ❌ DISABLED!\n# assert p99 < THRESHOLD_P99_MS   ❌ DISABLED!"
    r2.font.name = "Courier New"
    r2.font.size = Pt(16)
    r2.font.color.rgb = RGBColor(220, 53, 69)  # Red
    
    # Add link below code
    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(0.5))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    r_link = p2.add_run()
    r_link.text = "🔗 [View this code in editor →]"
    r_link.font.name = "Arial"
    r_link.font.size = Pt(18)
    r_link.font.color.rgb = RGBColor(0, 102, 204)
    r_link.hyperlink.address = create_code_link(
        "tests/integration/performance/test_performance_high_priority.py", "157"
    )
    
    # Impact
    tb3 = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1.5))
    tf3 = tb3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Impact:\n• 28 performance tests can't fail on poor performance\n• Only log warnings instead of blocking bad code"
    p3.font.name = "Arial"
    p3.font.size = Pt(18)


def add_code_evidence_2(prs):
    """Slide 4: Hardcoded 50% - WITH LINK"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Evidence from Code #2: Hardcoded 50%"
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.5))
    tf = tb.text_frame
    
    p1 = tf.paragraphs[0]
    r = p1.add_run()
    r.text = "# src/utils/validators.py:395\n\n"
    r.font.name = "Courier New"
    r.font.size = Pt(16)
    r.font.color.rgb = RGBColor(100, 100, 100)
    
    r2 = p1.add_run()
    r2.text = "def validate_roi_change_safety(\n    max_change_percent: float = 50.0  # ❌ NEVER CONFIRMED!\n):"
    r2.font.name = "Courier New"
    r2.font.size = Pt(16)
    r2.font.color.rgb = RGBColor(220, 53, 69)
    
    # Link
    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(0.5))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    r_link = p2.add_run()
    r_link.text = "🔗 [View this code in editor →]"
    r_link.font.name = "Arial"
    r_link.font.size = Pt(18)
    r_link.font.color.rgb = RGBColor(0, 102, 204)
    r_link.hyperlink.address = create_code_link("src/utils/validators.py", "395")
    
    # Impact
    tb3 = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(2))
    tf3 = tb3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Impact:\n• 6 ROI tests depend on unconfirmed value\n• Could be blocking legitimate use cases\n• Nobody knows if 50% is correct!"
    p3.font.name = "Arial"
    p3.font.size = Pt(18)


def add_top_7_critical_gaps(prs):
    """Slide 5: Top 7 Critical Gaps - WITH LINKS"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Top 7 Critical Gaps"
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    issues = [
        ("#1: Performance assertions disabled", "28 tests", 
         "test_performance_high_priority.py", "146"),
        
        ("#2: ROI 50% hardcoded", "6 tests",
         "validators.py", "395"),
        
        ("#3: NFFT validation too permissive", "6 tests",
         "validators.py", "194"),
        
        ("#4: Frequency range no maximum", "16 tests",
         "focus_server_models.py", "46"),
        
        ("#5: Sensor range no min/max", "15 tests",
         "validators.py", "116"),
        
        ("#6: API response time arbitrary", "3 tests",
         "test_api_endpoints_high_priority.py", "140"),
        
        ("#7: Config validation no assertions", "8 tests",
         "test_config_validation_high_priority.py", "475"),
    ]
    
    file_map = {
        "test_performance_high_priority.py": "tests/integration/performance/test_performance_high_priority.py",
        "validators.py": "src/utils/validators.py",
        "focus_server_models.py": "src/models/focus_server_models.py",
        "test_api_endpoints_high_priority.py": "tests/integration/api/test_api_endpoints_high_priority.py",
        "test_config_validation_high_priority.py": "tests/integration/api/test_config_validation_high_priority.py",
    }
    
    for i, (issue, count, file, line) in enumerate(issues):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        
        # Issue text
        r1 = p.add_run()
        r1.text = f"{issue} ({count})  "
        r1.font.name = "Arial"
        r1.font.size = Pt(16)
        
        # Link
        full_path = file_map.get(file, f"src/{file}")
        r2 = p.add_run()
        r2.text = "[view →]"
        r2.font.name = "Arial"
        r2.font.size = Pt(14)
        r2.font.color.rgb = RGBColor(0, 102, 204)
        r2.hyperlink.address = create_code_link(full_path, line)


def add_issue_details_slides(prs):
    """Slides 6-12: Issue Details - WITH LINKS WHERE RELEVANT"""
    
    # Slide 6: Performance SLAs
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Issue #1: Performance SLAs"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "What's Missing:\n• P95/P99 latency thresholds for each endpoint\n• Max error rate before failure\n• Different thresholds for live vs historic\n\nCurrently: Using guesses (500ms P95, 1000ms P99)"
    p.font.size = Pt(20)
    
    # Slide 7: ROI Change Limit
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Issue #2: ROI Change Limit"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Questions:\n• Is 50% correct?\n• Should it be 30%? 70%?\n• Is there a cooldown period?\n• Different limits for live vs historic?"
    p.font.size = Pt(22)
    
    # Slide 8: NFFT - WITH LINK
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Issue #3: NFFT Validation"
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Code only warns, never rejects:\nif not is_power_of_2:\n    warnings.warn(...)  # ⚠️\nreturn True  # ✅ Always passes!"
    r.font.name = "Courier New"
    r.font.size = Pt(18)
    
    # Link for NFFT
    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(0.5))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    r_link = p2.add_run()
    r_link.text = "🔗 [View validate_nfft_value in validators.py →]"
    r_link.font.size = Pt(16)
    r_link.font.color.rgb = RGBColor(0, 102, 204)
    r_link.hyperlink.address = create_code_link("src/utils/validators.py", "194")
    
    # Slides 9-12: Other issues (simpler, text-based)
    issues_simple = [
        ("Issue #4: Frequency Range", "Need absolute max/min frequency limits\nCurrently: Accepts any positive values"),
        ("Issue #5: Sensor Range", "Need min/max ROI size constraints\nCurrently: Could allow 1 sensor or all 2222"),
        ("Issue #6: API Response Times", "Need official SLAs for each endpoint\nCurrently: Arbitrary 1000ms timeout"),
        ("Issue #7: Config Validation", "Multiple edge cases with TODO/no assertions\nExample: min==max behavior undefined"),
    ]
    
    for issue_title, content in issues_simple:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = issue_title
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = content
        p.font.size = Pt(24)


def add_summary_slides(prs):
    """Slides 13-20: Cost, Solution, Timeline, Q&A, etc."""
    
    # Slide 13: Cost of Inaction
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Cost of Inaction"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "What happens if we don't define specs:\n• Tests exist but don't catch issues\n• False positives waste investigation time\n• False negatives miss real bugs\n• Can't confidently release\n• Automation loses credibility"
    p.font.size = Pt(22)
    
    # Slide 14: Proposed Solution
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Proposed Solution"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Three Steps:\n1. Schedule 2-3 hour specs meeting\n2. Define TOP 7 critical specs\n3. Update code + enable assertions\n\nTimeline: 1-2 weeks for implementation"
    p.font.size = Pt(24)
    
    # Slide 15: What We Need from Meeting
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "What We Need from Meeting"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Decisions on:\n• Performance thresholds (P95/P99)\n• ROI change limits\n• NFFT validation rules\n• Frequency/sensor range limits\n• API timeout values\n• Edge case behaviors"
    p.font.size = Pt(22)
    
    # Slides 16-19: Timeline, Expected Outcomes, Q&A Prep, Call to Action
    for slide_title, content in [
        ("Implementation Timeline", "Week 1: Update settings.py\nWeek 2: Enable assertions\nWeek 2: Re-run affected tests\nWeek 3: Document in Xray"),
        ("Expected Outcomes", "✅ Clear pass/fail criteria\n✅ Reliable automation\n✅ Confident releases\n✅ Reduced false positives"),
        ("Q&A Preparation", "Ready to answer:\n• Why now? (Blocking test development)\n• Can we estimate? (No, too risky)\n• How long? (2-3 hours meeting)"),
        ("Call to Action", "Let's schedule the specs meeting\nEvery day we wait = unreliable testing\n\nProposed: This week or next"),
    ]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = slide_title
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = content
        p.font.size = Pt(24)
    
    # Slide 20: Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Summary"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "📊 82+ tests affected\n🔗 Evidence in production code\n⏰ 2-3 hour meeting needed\n✅ Clear path forward\n\nLet's schedule it today!"
    p.font.size = Pt(28)


def main():
    """Generate full 20-slide presentation."""
    print("=" * 70)
    print("Generating FULL 20-Slide Presentation with Code Links")
    print("=" * 70)
    print()
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    print("[1/20] Creating title slide...")
    add_title_slide(prs)
    
    print("[2/20] Creating executive summary...")
    add_executive_summary(prs)
    
    print("[3/20] Creating code evidence #1 (with link)...")
    add_code_evidence_1(prs)
    
    print("[4/20] Creating code evidence #2 (with link)...")
    add_code_evidence_2(prs)
    
    print("[5/20] Creating Top 7 gaps (with links)...")
    add_top_7_critical_gaps(prs)
    
    print("[6-12/20] Creating issue detail slides...")
    add_issue_details_slides(prs)
    
    print("[13-20/20] Creating summary slides...")
    add_summary_slides(prs)
    
    prs.save(DECK_PATH)
    
    print()
    print("=" * 70)
    print("COMPLETE!")
    print("=" * 70)
    print()
    print(f"File created: {DECK_PATH}")
    print()
    print("What's included:")
    print("  - 20 slides total")
    print("  - Code links in slides: 3, 4, 5, 8")
    print("  - All Top 7 issues have clickable links")
    print("  - Professional formatting")
    print()
    print("Test it:")
    print(f"  start {DECK_PATH}")
    print()


if __name__ == "__main__":
    main()

